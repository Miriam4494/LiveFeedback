import os
import boto3
import tempfile
import mimetypes
from dotenv import load_dotenv
import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
from sentence_transformers import SentenceTransformer
import pinecone
import uuid
from urllib.parse import urlparse
from pptx import Presentation
import zipfile



load_dotenv()  # טוען את משתני הסביבה מהקובץ .env

openai_api_key = os.getenv("OPENAI_API_KEY")
pinecone_api_key = os.getenv("PINECONE_API_KEY")

from urllib.parse import unquote
import openai
from pinecone import Pinecone, ServerlessSpec
openai.api_key = openai_api_key
# Initialize Pinecone instance
pinecone = Pinecone(
    api_key=pinecone_api_key,
    ssl_verify=False
)
spec = ServerlessSpec(
    cloud="aws",
    region="us-east-1"
)
pinecone_index_name = "user-files"

# מודל לאמבדינג
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

# הורדת קובץ מ-S3
import requests
import tempfile

def create_presigned_url(bucket_name, object_key, expiration=3600):
    s3_client = boto3.client(
        's3',
        aws_access_key_id=os.getenv("AWS_ACCES_KEY_ID"),
        aws_secret_access_key=os.getenv("AWS_SECRET_ACCES_KEY"),
        )
    try:
        response = s3_client.generate_presigned_url('get_object',
                                                    Params={'Bucket': bucket_name, 'Key': object_key},
                                                    ExpiresIn=expiration)
    except Exception as e:
        print(f"❗ Error generating presigned URL: {e}")
        raise
    return response


def download_s3_file(s3_url):
    print(f"📥 Downloading from URL: {s3_url}")
    try:
        response = requests.get(s3_url, stream=True)
        response.raise_for_status()
        content_type = response.headers.get('Content-Type', '')
        extension = mimetypes.guess_extension(content_type) or ''
        file_name = "unknown_file"

        if not extension:
            parsed_url = urlparse(s3_url)
            file_name = os.path.basename(parsed_url.path)
            extension = os.path.splitext(file_name)[1]  # קבלת הסיומת מהשם
        print(f"📄 שם הקובץ מה-URL: {file_name}")
        print(f"📄 סיומת הקובץ: {extension}")
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=extension)

        with open(temp_file.name, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        print(f"📂 קובץ זמני נשמר בנתיב: {temp_file.name}")
        return temp_file.name
    except Exception as e:
        print(f"❗ Error downloading file: {e}")
        raise

from PIL import Image

def validate_and_convert_image(file_path):
    """
    המרת תמונה לפורמט PNG אם נדרש
    """
    try:
        img = Image.open(file_path)
        converted_path = file_path + ".png"
        img.save(converted_path, format="PNG")
        return converted_path
    except Exception as e:
        print(f"❗ שגיאה בהמרת התמונה: {e}")
        return None


def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower().replace('.', '')  # תמיד נשתמש בסיומת כמו 'txt'
    print(f"📂 נתיב הקובץ: {file_path}")
    print(f"📄 סיומת הקובץ: {ext}")
    text = ""

    try:
        if ext in ['txt', 'md', 'json', 'csv', 'py', 'js', 'html', 'css']:
            # קריאת קבצי טקסט
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()

        elif ext == 'pdf':
            # קריאת קבצי PDF
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""

        elif ext == 'docx':
            # קריאת קבצי Word
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"

        elif ext in ['jpeg', 'jpg', 'png']:
             mime_type = f"image/{ext if ext != 'jpg' else 'jpeg'}"
             text = transcribe_image_with_ai(file_path, mime_type)
        elif ext in ['mp4', 'mov', 'avi', 'mkv']:
            # תמלול וידאו באמצעות AI
            text = transcribe_video_with_ai(file_path)
        elif ext == 'zip':
            text = extract_text_from_zip(file_path)
        elif ext == 'pptx':
            text = extract_text_from_pptx(file_path)
        elif ext in ['mp3', 'wav', 'ogg']:
            text = transcribe_audio_with_ai(file_path)
        else:
            print(f"⚠️ סיומת לא נתמכת: {ext}")
            text = "סוג קובץ לא נתמך לקריאה ישירה."
    except Exception as e:
        print(f"❗ שגיאה בקריאת הקובץ: {e}")
        text = ""

    return text

import base64

def extract_text_from_zip(file_path):
    """
    חילוץ טקסט מכל הקבצים הזמינים ב-ZIP (תוך התעלמות מקבצים בינאריים)
    """
    text = ""
    try:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            for file_info in zip_ref.infolist():
                if not file_info.filename.endswith(('.jpg', '.png', '.mp4', '.mov', '.avi', '.mp3', '.wav')):
                    with zip_ref.open(file_info.filename) as file:
                        try:
                            content = file.read().decode('utf-8', errors='ignore')
                            text += f"\n\n--- {file_info.filename} ---\n\n{content}"
                        except Exception as e:
                            print(f"שגיאה בקריאת {file_info.filename}: {e}")
    except Exception as e:
        print(f"שגיאה בפענוח קובץ ZIP: {e}")
    return text



def extract_text_from_pptx(file_path):
    """
    שליפת טקסט ממצגת PowerPoint
    """
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"שגיאה בקריאת PPTX: {e}")
    return text

def transcribe_audio_with_ai(file_path):
    """
    תמלול קובץ אודיו באמצעות Whisper
    """
    try:
        with open(file_path, "rb") as audio_file:
            transcript = openai.audio.transcriptions.create(
                model="whisper-1",
                file=audio_file
            )
        return transcript.text
    except Exception as e:
        print(f"שגיאה בתמלול האודיו: {e}")
        return f"שגיאה בתמלול האודיו: {str(e)}"


def transcribe_image_with_ai(file_path,mime_type):
    """
    שליחת תמונה ל-OpenAI עם פרומפט מותאם לקבלת תיאור של התמונה
    """
    print(f"📤 שליחת תמונה ל-OpenAI לתיאור: {file_path}")
    try:
        with open(file_path, "rb") as img_file:
            b64 = base64.b64encode(img_file.read()).decode("utf-8")
        data_url = f"data:{mime_type};base64,{b64}"

        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "user", "content": [
                    {"type": "text", "text": " תאר לי את מה שרואים בתמונה הזו כולל מילים שרשומות בה, צבעים, פרטים וכו'."},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ]}
            ]
        )
        return response.choices[0].message.content

    except Exception as e:
        print(f"❗ שגיאה בשליחת התמונה ל-OpenAI: {e}")
        return f"❗ שגיאה בניתוח התמונה עם OpenAI: {str(e)}"
   
def transcribe_video_with_ai(file_path):
    """
    תמלול וידאו לקובץ טקסט בעזרת OpenAI Whisper
    """
    print(f"📤 שליחת תמונה ל-OpenAI לתיאור: {file_path}")
    try:
        with open(file_path, "rb") as video_file:
            transcript = openai.audio.transcriptions.create(
                model="whisper-1",
                file=video_file
            )
        return transcript.text

    except Exception as e:
        print(f"❗ שגיאה בתמלול הוידאו: {e}")
        return f"❗ שגיאה בתמלול הוידאו: {str(e)}"

  
# פיצול טקסט לפסקאות קטנות
def split_text(text, max_chunk_size=500):
    paragraphs = text.split("\n\n")
    chunks = []
    current_chunk = ""

    for para in paragraphs:
        if len(current_chunk) + len(para) <= max_chunk_size:
            current_chunk += para + "\n\n"
        else:
            chunks.append(current_chunk.strip())
            current_chunk = para + "\n\n"

    if current_chunk:
        chunks.append(current_chunk.strip())

    return chunks

# פונקציה עיקרית: קבלת קובץ, אינדוקס ושמירה ב-Pinecone
def index_s3_file_for_user(s3_url: str, user_id: str,file_id: str):
    # שלב 1: הורדה
    print(f"📥 הורדת הקובץ מ-S3: {s3_url}")
    parsed_url = urlparse(s3_url)
    bucket_name = parsed_url.netloc.split('.')[0]  
    # object_key = parsed_url.path.lstrip('/')  
    object_key = unquote(parsed_url.path.lstrip('/'))  # ✅ חשוב


    print("📂 object_key:", object_key)

    # יצירת Presigned URL
    presigned_url = create_presigned_url(bucket_name, object_key)

    # הורדת הקובץ
    local_file_path = download_s3_file(presigned_url)    
    # שלב 2: קריאת התוכן
    text = extract_text(local_file_path)
    if not text.strip():
        print(f"⚠️ הקובץ ריק או לא נתמך: {s3_url}")
        os.remove(local_file_path)
        return

    # שלב 3: חלוקה לחלקים
    text_chunks = split_text(text)

    # שלב 4: יצירת אמבדינגים
    embeddings = embedding_model.encode(text_chunks)

    # שלב 5: שליחה ל-Pinecone
    if pinecone_index_name not in pinecone.list_indexes().names():
        pinecone.create_index(
            name=pinecone_index_name,
            dimension=len(embeddings[0]),
            metric="cosine",
            spec=spec
        )

    index = pinecone.Index(pinecone_index_name)

    vectors = []
    for idx, (chunk, embedding) in enumerate(zip(text_chunks, embeddings)):
        vectors.append({
            "id": f"{user_id}_{uuid.uuid4().hex}",
            "values": embedding.tolist(),
            "metadata": {
                "user_id": user_id,
                "file_id": file_id,
                "text": chunk
            }
        })

    index.upsert(vectors)

    print(f"✔️ {len(vectors)} קטעים הוכנסו ל-Pinecone תחת משתמש {user_id} מהקובץ {file_id}")

    # ניקוי קובץ זמני
    os.remove(local_file_path)


def describe_file_from_url(file_url):
    # שלב 1: הורדת הקובץ
    response = requests.get(file_url)
    if response.status_code != 200:
        return "Failed to download the file."

    # שלב 2: זיהוי סוג MIME
    content_type = response.headers.get('Content-Type')
    ext = mimetypes.guess_extension(content_type)

    # שלב 3: שמירה זמנית
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp_file:
        tmp_file.write(response.content)
        tmp_file_path = tmp_file.name

    # שלב 4: הכנה לשליחה ל-AI
    if 'image' in content_type:
        # תמונה - שליחה ל-GPT-4-Vision
        with open(tmp_file_path, "rb") as image_file:
            result = openai.ChatCompletion.create(
                model="gpt-4-vision-preview",
                messages=[
                    {"role": "user", "content": [
                        {"type": "text", "text": "Please describe the content of this image."},
                        {"type": "image_url", "image_url": {"url": f"data:{content_type};base64,{image_file.read().encode('base64')}"}}  # שימי לב – רק אם יש תמיכה
                    ]}
                ],
                max_tokens=300
            )
            return result.choices[0].message['content']

    elif 'pdf' in content_type or 'text' in content_type:
        # טקסטים – קריאה ושליחה ל-GPT
        text = response.content.decode(errors='ignore')[:2000]  # ניקח רק תחילת הקובץ כדי לא לעבור מגבלת טוקנים
        result = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[{"role": "user", "content": f"Please describe the content of this file:\n{text}"}],
            max_tokens=300
        )
        return result.choices[0].message['content']

    elif 'audio' in content_type or 'video' in content_type:
        # שליחה ל-Whisper (לדוגמה, רק אודיו)
        with open(tmp_file_path, "rb") as media_file:
            transcript = openai.Audio.transcribe("whisper-1", media_file)
            return f"Transcription of the file:\n{transcript['text']}"

    else:
        return "Unsupported file type for automatic description."

def query_user_files(user_id: str, query: str, score_threshold: float = 0.8):
    """
    פונקציה לחיפוש קבצים לפי שאילתה ו-user_id
    """
    print(f"🔍 חיפוש קבצים עבור user_id: {user_id} עם שאילתה: {query}")
    try:
        # יצירת אמבדינג לשאילתה
        query_embedding = embedding_model.encode([query])[0]

        # בדיקת קיום האינדקס
        if pinecone_index_name not in pinecone.list_indexes().names():
            raise ValueError(f"⚠️ אינדקס {pinecone_index_name} לא קיים ב-Pinecone.")

        # חיפוש באינדקס
        index = pinecone.Index(pinecone_index_name)
        results = index.query(
            vector=query_embedding.tolist(),
            top_k=100,
            include_metadata=True,
            filter={"user_id": user_id}  # סינון לפי user_id
        )

        # עיבוד התוצאות
        query_results = []
        for match in results["matches"]:
            if match["score"] >= score_threshold:  # סינון לפי הסף
                query_results.append({
                    "file_id": match["metadata"]["file_id"],
                    "text_snippet": match["metadata"]["text"],
                    "score": match["score"]
                })

        return query_results
    except Exception as e:
        print(f"❗ שגיאה בחיפוש קבצים: {e}")
        raise


from fastapi import FastAPI
from pydantic import BaseModel
from typing import List
import uvicorn



app = FastAPI()

# ---- מודלים ל-Request ----

class IndexFileRequest(BaseModel):
    s3_url: str
    user_id: int
    file_id: int

class QueryFilesRequest(BaseModel):
    user_id: int
    query: str
    score_threshold: float = 0.1  # סף ציון ברירת מחדל
class QueryResult(BaseModel):
    file_id: int
    text_snippet: str
    score: float

# ---- ENDPOINTS ----

@app.post("/index-file")
def index_file(req: IndexFileRequest):
    try:
        print("req: ",req)
        index_s3_file_for_user(req.s3_url, req.user_id, req.file_id)
        return {"status": "success", "message": "File indexed successfully."}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/query-files", response_model=List[QueryResult])
def query_files(req: QueryFilesRequest):
    try:
        results = query_user_files(req.user_id, req.query, req.score_threshold)
        return results
    except Exception as e:
        return [{"file_id": "", "text_snippet": f"Error: {str(e)}", "score": 0.0}]

# ---- הרצה ----

import nest_asyncio
nest_asyncio.apply()

uvicorn.run(app, host="0.0.0.0", port=8000)
